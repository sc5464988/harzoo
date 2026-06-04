"""Document reader tool supporting common office/text formats."""


from __future__ import annotations

import csv
import json
from pathlib import Path
from typing import Any
from xml.etree import ElementTree

import yaml

from harzoo.agent.kernel.tool import Tool, ToolResult

TOOL_VERSION = "2026-05-22"


def _resolve_path(path: str) -> Path:
    p = Path(path).expanduser()
    if p.is_absolute():
        return p.resolve()
    return (Path.cwd() / p).resolve()


def _truncate_text(text: str, max_chars: int) -> tuple[str, bool]:
    if len(text) <= max_chars:
        return text, False
    return text[:max_chars], True


class DocumentReadTool(Tool):
    """文档读取工具：将多种文档格式统一转换为文本/结构化输出。"""

    name = "DocumentRead"
    description = "Read common documents (pdf/docx/xlsx/xls/pptx/csv/txt/json/yaml/xml/html) as text, markdown, or structured output."
    parameters = {
        "type": "object",
        "properties": {
            "file_path": {"type": "string", "description": "Path to the document"},
            "mode": {
                "type": "string",
                "enum": ["text", "markdown", "structured"],
                "description": "Output mode. Defaults to text.",
                "default": "text",
            },
            "page": {"type": "integer", "description": "Optional 1-based page/slide index for pdf/pptx"},
            "sheet": {"type": "string", "description": "Optional sheet name for xlsx/xls"},
            "max_chars": {
                "type": "integer",
                "description": "Max characters for content field",
                "default": 50000,
                "minimum": 1000,
                "maximum": 300000,
            },
            "include_meta": {"type": "boolean", "description": "Include metadata in output", "default": True},
        },
        "required": ["file_path"],
    }
    risk_level = "low"
    side_effect = False
    idempotent = True

    _TEXT_EXTENSIONS = {".txt", ".md", ".json", ".yaml", ".yml", ".xml", ".html", ".htm", ".csv"}
    _OFFICE_EXTENSIONS = {".pdf", ".docx", ".xlsx", ".xls", ".pptx"}
    _SUPPORTED_EXTENSIONS = _TEXT_EXTENSIONS | _OFFICE_EXTENSIONS

    def execute(
        self,
        file_path: str,
        mode: str = "text",
        page: int | None = None,
        sheet: str | None = None,
        max_chars: int = 50000,
        include_meta: bool = True,
        **_: Any,
    ) -> ToolResult:
        selected_mode = str(mode or "text").strip().lower()
        if selected_mode not in {"text", "markdown", "structured"}:
            return ToolResult.failure("mode must be one of: text, markdown, structured", code="INVALID_ARGUMENTS")
        if max_chars < 1000:
            return ToolResult.failure("max_chars must be >= 1000", code="INVALID_ARGUMENTS")
        if page is not None and int(page) < 1:
            return ToolResult.failure("page must be >= 1", code="INVALID_ARGUMENTS")

        p = _resolve_path(file_path)
        if not p.exists() or not p.is_file():
            return ToolResult.failure(f"File not found: {file_path}", code="FILE_NOT_FOUND")

        ext = p.suffix.lower()
        if ext not in self._SUPPORTED_EXTENSIONS:
            return ToolResult.failure(f"Unsupported document format: {ext or '(no extension)'}", code="UNSUPPORTED_FORMAT")

        try:
            doc = self._read_document(p, ext=ext, page=page, sheet=sheet)
            data = self._render(doc, mode=selected_mode, max_chars=max_chars, include_meta=bool(include_meta))
            return ToolResult.success(data)
        except ToolRuntimeError as exc:
            return ToolResult.failure(str(exc), code=exc.code)
        except Exception as exc:  # noqa: BLE001
            return ToolResult.failure(f"{type(exc).__name__}: {exc}", code="TOOL_EXCEPTION")

    def _read_document(self, p: Path, *, ext: str, page: int | None, sheet: str | None) -> dict[str, Any]:
        if ext in {".txt", ".md"}:
            text = p.read_text(encoding="utf-8", errors="replace")
            return {"format": ext[1:], "blocks": [{"type": "text", "text": text}], "tables": [], "warnings": [], "meta": {}}
        if ext == ".json":
            obj = json.loads(p.read_text(encoding="utf-8", errors="replace"))
            text = json.dumps(obj, ensure_ascii=False, indent=2)
            return {"format": "json", "blocks": [{"type": "text", "text": text}], "tables": [], "warnings": [], "meta": {}}
        if ext in {".yaml", ".yml"}:
            obj = yaml.safe_load(p.read_text(encoding="utf-8", errors="replace"))
            text = yaml.safe_dump(obj, allow_unicode=True, sort_keys=False)
            return {"format": "yaml", "blocks": [{"type": "text", "text": text}], "tables": [], "warnings": [], "meta": {}}
        if ext in {".xml", ".html", ".htm"}:
            text = self._read_markup(p)
            fmt = "html" if ext in {".html", ".htm"} else "xml"
            return {"format": fmt, "blocks": [{"type": "text", "text": text}], "tables": [], "warnings": [], "meta": {}}
        if ext == ".csv":
            return self._read_csv(p)
        if ext == ".pdf":
            return self._read_pdf(p, page=page)
        if ext == ".docx":
            return self._read_docx(p)
        if ext == ".xlsx":
            return self._read_xlsx(p, sheet=sheet)
        if ext == ".xls":
            return self._read_xls(p, sheet=sheet)
        if ext == ".pptx":
            return self._read_pptx(p, page=page)
        raise ToolRuntimeError(f"Unsupported document format: {ext}", code="UNSUPPORTED_FORMAT")

    def _read_markup(self, p: Path) -> str:
        try:
            root = ElementTree.fromstring(p.read_text(encoding="utf-8", errors="replace"))
            text = "".join(root.itertext())
            return "\n".join(line.strip() for line in text.splitlines() if line.strip())
        except ElementTree.ParseError:
            return p.read_text(encoding="utf-8", errors="replace")

    def _read_csv(self, p: Path) -> dict[str, Any]:
        rows: list[list[str]] = []
        with p.open("r", encoding="utf-8", errors="replace", newline="") as fh:
            reader = csv.reader(fh)
            for row in reader:
                rows.append([str(cell) for cell in row])
        table_text = "\n".join(", ".join(r) for r in rows)
        return {
            "format": "csv",
            "blocks": [{"type": "text", "text": table_text}],
            "tables": [{"source": p.name, "rows": rows}],
            "warnings": [],
            "meta": {"row_count": len(rows)},
        }

    def _read_pdf(self, p: Path, *, page: int | None) -> dict[str, Any]:
        try:
            from pypdf import PdfReader  # pyright: ignore[reportMissingImports]
        except Exception as exc:  # noqa: BLE001
            raise ToolRuntimeError(f"PDF dependency missing: {exc}", code="PARSE_FAILED") from exc
        reader = PdfReader(str(p))
        pages = reader.pages
        indices = [page - 1] if page is not None else list(range(len(pages)))
        blocks: list[dict[str, Any]] = []
        tables: list[dict[str, Any]] = []
        warnings: list[str] = []
        plumber_pdf = None
        try:
            import pdfplumber  # type: ignore

            plumber_pdf = pdfplumber.open(str(p))
        except Exception:
            warnings.append("pdfplumber unavailable; table extraction skipped.")
        for idx in indices:
            if idx < 0 or idx >= len(pages):
                raise ToolRuntimeError(f"Requested page out of range: {idx + 1}", code="INVALID_ARGUMENTS")
            text = pages[idx].extract_text() or ""
            if not text.strip():
                warnings.append(f"No extractable text on page {idx + 1}.")
            blocks.append({"type": "page", "page": idx + 1, "text": text.strip()})
            if plumber_pdf is not None and idx < len(plumber_pdf.pages):
                extracted_tables = plumber_pdf.pages[idx].extract_tables() or []
                for table_i, raw_rows in enumerate(extracted_tables, start=1):
                    rows = [
                        [("" if cell is None else str(cell).strip()) for cell in row]
                        for row in raw_rows
                        if row is not None
                    ]
                    if rows:
                        tables.append({"source": f"page_{idx + 1}_table_{table_i}", "page": idx + 1, "rows": rows})
        if plumber_pdf is not None:
            plumber_pdf.close()
        return {"format": "pdf", "blocks": blocks, "tables": tables, "warnings": warnings, "meta": {"pages": len(pages), "tables": len(tables)}}

    def _read_docx(self, p: Path) -> dict[str, Any]:
        try:
            from docx import Document  # pyright: ignore[reportMissingImports]
        except Exception as exc:  # noqa: BLE001
            raise ToolRuntimeError(f"DOCX dependency missing: {exc}", code="PARSE_FAILED") from exc
        doc = Document(str(p))
        blocks: list[dict[str, Any]] = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                blocks.append({"type": "paragraph", "text": text})
        tables = []
        for i, table in enumerate(doc.tables, start=1):
            rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
            tables.append({"source": f"table_{i}", "rows": rows})
        return {"format": "docx", "blocks": blocks, "tables": tables, "warnings": [], "meta": {"paragraphs": len(blocks), "tables": len(tables)}}

    def _read_xlsx(self, p: Path, *, sheet: str | None) -> dict[str, Any]:
        try:
            from openpyxl import load_workbook  # pyright: ignore[reportMissingImports,reportMissingModuleSource]
        except Exception as exc:  # noqa: BLE001
            raise ToolRuntimeError(f"XLSX dependency missing: {exc}", code="PARSE_FAILED") from exc
        wb = load_workbook(str(p), data_only=True, read_only=True)
        sheet_names = wb.sheetnames
        selected = [sheet] if sheet else sheet_names
        tables = []
        blocks = []
        for name in selected:
            if name not in wb:
                raise ToolRuntimeError(f"Sheet not found: {name}", code="INVALID_ARGUMENTS")
            ws = wb[name]
            rows = [[("" if cell is None else str(cell)) for cell in row] for row in ws.iter_rows(values_only=True)]
            tables.append({"source": name, "rows": rows})
            blocks.append({"type": "sheet", "sheet": name, "text": "\n".join(", ".join(r) for r in rows)})
        return {"format": "xlsx", "blocks": blocks, "tables": tables, "warnings": [], "meta": {"sheets": sheet_names}}

    def _read_xls(self, p: Path, *, sheet: str | None) -> dict[str, Any]:
        try:
            import xlrd  # pyright: ignore[reportMissingImports,reportMissingModuleSource]
        except Exception as exc:  # noqa: BLE001
            raise ToolRuntimeError(f"XLS dependency missing: {exc}", code="PARSE_FAILED") from exc
        book = xlrd.open_workbook(str(p))
        sheet_names = book.sheet_names()
        selected = [sheet] if sheet else sheet_names
        tables = []
        blocks = []
        for name in selected:
            if name not in sheet_names:
                raise ToolRuntimeError(f"Sheet not found: {name}", code="INVALID_ARGUMENTS")
            ws = book.sheet_by_name(name)
            rows: list[list[str]] = []
            for row_idx in range(ws.nrows):
                rows.append([str(ws.cell_value(row_idx, col_idx)) for col_idx in range(ws.ncols)])
            tables.append({"source": name, "rows": rows})
            blocks.append({"type": "sheet", "sheet": name, "text": "\n".join(", ".join(r) for r in rows)})
        return {"format": "xls", "blocks": blocks, "tables": tables, "warnings": [], "meta": {"sheets": sheet_names}}

    def _read_pptx(self, p: Path, *, page: int | None) -> dict[str, Any]:
        try:
            from pptx import Presentation  # pyright: ignore[reportMissingImports]
        except Exception as exc:  # noqa: BLE001
            raise ToolRuntimeError(f"PPTX dependency missing: {exc}", code="PARSE_FAILED") from exc
        prs = Presentation(str(p))
        slides = prs.slides
        indices = [page - 1] if page is not None else list(range(len(slides)))
        blocks: list[dict[str, Any]] = []
        for idx in indices:
            if idx < 0 or idx >= len(slides):
                raise ToolRuntimeError(f"Requested slide out of range: {idx + 1}", code="INVALID_ARGUMENTS")
            slide = slides[idx]
            parts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and str(getattr(shape, "text", "")).strip():
                    parts.append(str(shape.text).strip())
            note_text = ""
            if slide.has_notes_slide and slide.notes_slide and slide.notes_slide.notes_text_frame:
                note_text = str(slide.notes_slide.notes_text_frame.text or "").strip()
            block: dict[str, Any] = {"type": "slide", "slide": idx + 1, "text": "\n".join(parts)}
            if note_text:
                block["notes"] = note_text
            blocks.append(block)
        return {"format": "pptx", "blocks": blocks, "tables": [], "warnings": [], "meta": {"slides": len(slides)}}

    def _render(self, doc: dict[str, Any], *, mode: str, max_chars: int, include_meta: bool) -> dict[str, Any]:
        sections = doc.get("blocks", [])
        tables = doc.get("tables", [])
        warnings = list(doc.get("warnings", []))

        lines: list[str] = []
        for b in sections:
            text = str(b.get("text") or "").strip()
            if text:
                lines.append(text)
            notes = str(b.get("notes") or "").strip()
            if notes:
                lines.append(notes)
        for table in tables:
            rows = table.get("rows") or []
            for row in rows:
                lines.append(", ".join(str(cell) for cell in row))

        joined = "\n\n".join(line for line in lines if line)
        content, truncated = _truncate_text(joined, max_chars)

        if mode == "structured":
            out: dict[str, Any] = {
                "content": content,
                "sections": sections,
                "tables": tables,
                "warnings": warnings,
            }
        elif mode in {"text", "markdown"}:
            out = {"content": content, "warnings": warnings}
        else:
            raise ToolRuntimeError(f"Unsupported mode: {mode}", code="INVALID_ARGUMENTS")

        if include_meta:
            meta = dict(doc.get("meta", {}))
            meta["format"] = doc.get("format")
            meta["truncated"] = truncated
            out["meta"] = meta
        return out


class ToolRuntimeError(Exception):
    """工具内部业务异常，携带可对外暴露的错误码。"""
    def __init__(self, message: str, *, code: str) -> None:
        super().__init__(message)
        self.code = code


TOOL = DocumentReadTool
